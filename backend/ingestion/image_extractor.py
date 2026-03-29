import fitz  # for pdf images
import os
import base64
from docx import Document  # for docx images
from pptx import Presentation # for pptx images
from config import UPLOAD_DIR

def extract_images(file_path: str, doc_id: str) -> list:
    ext = os.path.splitext(file_path)[-1].lower()
    
    image_dir = os.path.join(UPLOAD_DIR, doc_id, "images")
    os.makedirs(image_dir, exist_ok=True)

    if ext == ".pdf":
        return _extract_from_pdf(file_path, doc_id, image_dir)
    elif ext == ".docx":
        return _extract_from_docx(file_path, doc_id, image_dir)
    elif ext == ".pptx":
        return _extract_from_pptx(file_path, doc_id, image_dir)
    else:
        return []  # CSV, XLSX — no images

def _save_image(image_bytes: bytes, image_path: str) -> str:
    with open(image_path, "wb") as f:
        f.write(image_bytes)
    return image_path

def _extract_from_pdf(file_path: str, doc_id: str, image_dir: str) -> list:
    doc = fitz.open(file_path)
    images = []

    for page_num, page in enumerate(doc):
        for img_index, img in enumerate(page.get_images(full=True)):  # list of image references (xref).
            xref = img[0]
            base_image = doc.extract_image(xref)  #extracts bytes and format
            image_bytes = base_image["image"]
            ext = base_image["ext"]

            image_filename = f"page{page_num + 1}_img{img_index}.{ext}"
            image_path = os.path.join(image_dir, image_filename)
            _save_image(image_bytes, image_path)

            images.append({
                "doc_id": doc_id,
                "doc_name": os.path.basename(file_path),
                "page_num": page_num + 1,
                "image_index": img_index,
                "image_path": image_path,
                "image_b64": base64.b64encode(image_bytes).decode("utf-8"), # can be directly feed into json and to vision model.
                "source": "pdf"
            })

    doc.close()
    return images

def _extract_from_docx(file_path: str, doc_id: str, image_dir: str) -> list:
    doc = Document(file_path)
    images = []
    img_index = 0

    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            image_bytes = rel.target_part.blob
            ext = rel.target_part.content_type.split("/")[-1]

            image_filename = f"section0_img{img_index}.{ext}"
            image_path = os.path.join(image_dir, image_filename)
            _save_image(image_bytes, image_path)

            images.append({
                "doc_id": doc_id,
                "doc_name": os.path.basename(file_path),
                "page_num": None,
                "section_index": 0,
                "image_index": img_index,
                "image_path": image_path,
                "image_b64": base64.b64encode(image_bytes).decode("utf-8"),
                "source": "docx"
            })
            img_index += 1

    return images

def _extract_from_pptx(file_path: str, doc_id: str, image_dir: str) -> list:
    prs = Presentation(file_path)
    images = []

    for slide_index, slide in enumerate(prs.slides):
        img_index = 0
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image_bytes = shape.image.blob
                ext = shape.image.ext

                image_filename = f"slide{slide_index + 1}_img{img_index}.{ext}"
                image_path = os.path.join(image_dir, image_filename)
                _save_image(image_bytes, image_path)

                images.append({
                    "doc_id": doc_id,
                    "doc_name": os.path.basename(file_path),
                    "page_num": None,
                    "slide_index": slide_index + 1,
                    "image_index": img_index,
                    "image_path": image_path,
                    "image_b64": base64.b64encode(image_bytes).decode("utf-8"),
                    "source": "pptx"
                })
                img_index += 1

    return images